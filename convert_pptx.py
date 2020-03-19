import os
import argparse

from pptx import Presentation

from converter import *

parser = argparse.ArgumentParser()
parser.add_argument('file')

if __name__ == "__main__":
    args = parser.parse_args()
    prs = Presentation(args.file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for i in range(len(paragraph.runs)):
                    paragraph.runs[i].text = CovertEncoding(text, _VNIWin, _Unicode)

    prs.save(file[:-5]+"_unicode.pptx")